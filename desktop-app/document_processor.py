"""
Document processor - extracts text from various file formats and images via OCR
"""
import os
from pathlib import Path
from typing import Optional, Dict, Any
import json
import logging

# Document libraries
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import pandas as pd

# OCR support
try:
    from ocr.cli import process_image, process_pdf
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Process and extract text from documents"""
    
    def __init__(self, database):
        """Initialize with database reference"""
        self.db = database
        
    def process_document(self, file_path: str) -> bool:
        """Process a document and add to database"""
        try:
            # Get file info
            filename = os.path.basename(file_path)
            file_size = os.path.getsize(file_path)
            file_type = os.path.splitext(filename)[1].lower()
            
            # Extract text
            extracted_text = self.extract_text_from_file(file_path)
            
            # Add to database
            doc_id = self.db.add_document(
                filename=filename,
                file_path=file_path,
                file_size=file_size,
                file_type=file_type,
                extracted_text=extracted_text,
                metadata=json.dumps({})
            )
            
            return doc_id != -1
            
        except Exception as e:
            print(f"Error processing {file_path}: {e}")
            return False
            
    def extract_text_from_file(self, file_path: str) -> str:
        """Extract text from various file formats"""
        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext == '.pdf':
                return self._extract_from_pdf(file_path)
            elif ext == '.txt' or ext == '.md':
                return self._extract_from_txt(file_path)
            elif ext == '.docx':
                return self._extract_from_docx(file_path)
            elif ext == '.pptx':
                return self._extract_from_pptx(file_path)
            elif ext == '.xlsx':
                return self._extract_from_xlsx(file_path)
            elif ext == '.csv':
                return self._extract_from_csv(file_path)
            else:
                return ""
                
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return f"Error extracting text: {str(e)}"
            
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF"""
        reader = PdfReader(file_path)
        text_parts = []
        
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
                
        return "\n\n".join(text_parts)
        
    def _extract_from_txt(self, file_path: str) -> str:
        """Extract text from TXT/MD files"""
        encodings = ['utf-8', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
                
        return "Unable to decode file"
        
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        doc = DocxDocument(file_path)
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
                
        return "\n\n".join(text_parts)
        
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX"""
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                    
            if slide_text:
                text_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
                
        return "\n\n".join(text_parts)
        
    def _extract_from_xlsx(self, file_path: str) -> str:
        """Extract text from XLSX"""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = [f"Sheet: {sheet_name}"]
            
            for row in ws.iter_rows(values_only=True, max_row=100):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    sheet_text.append(row_text)
                    
            text_parts.append("\n".join(sheet_text))
            
        return "\n\n".join(text_parts)
    
    def _extract_from_image_ocr(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from image using OCR.
        
        Returns:
            {
                "text": str,
                "confidence": float,
                "method": "ocr",
                "boxes": [...],
                "language": str
            }
        """
        if not HAS_OCR:
            logger.warning("OCR not available; skipping image processing")
            return {
                "text": "",
                "confidence": 0.0,
                "method": "ocr",
                "error": "OCR module not installed"
            }
        
        try:
            result = process_image(file_path)
            
            if "error" in result:
                logger.error(f"OCR error: {result['error']}")
                return result
            
            return {
                "text": result.get("text", ""),
                "confidence": result.get("confidence", 0.0),
                "method": "ocr",
                "lines": result.get("lines", []),
                "language": result.get("metadata", {}).get("language", "en")
            }
        except Exception as e:
            logger.error(f"Error processing image with OCR: {e}")
            return {
                "text": "",
                "confidence": 0.0,
                "method": "ocr",
                "error": str(e)
            }
    
    def _extract_from_pdf_ocr(self, file_path: str, dpi: int = 150) -> str:
        """
        Extract text from PDF using OCR.
        
        Args:
            file_path: Path to PDF
            dpi: Resolution for rendering
        
        Returns:
            Combined text from all pages
        """
        if not HAS_OCR:
            logger.warning("OCR not available; cannot process PDF")
            return ""
        
        try:
            results = process_pdf(file_path, dpi=dpi)
            text_parts = []
            
            for result in results:
                if "error" not in result:
                    page_num = result.get("page", 0)
                    text_parts.append(f"--- Page {page_num} ---\n{result.get('text', '')}")
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error processing PDF with OCR: {e}")
            return f"Error: {e}"
            
        return "\n\n".join(text_parts)
        
    def _extract_from_csv(self, file_path: str) -> str:
        """Extract text from CSV"""
        try:
            df = pd.read_csv(file_path, nrows=100)
            return df.to_string()
        except Exception as e:
            return f"Error reading CSV: {str(e)}"
