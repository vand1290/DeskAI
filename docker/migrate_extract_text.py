"""
Migration script to extract text from existing documents
"""
import psycopg2
from minio import Minio
import io
from pathlib import Path
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
import pandas as pd

def extract_text_from_file(file_bytes, filename):
    """Extract text content from various file types"""
    try:
        file_ext = Path(filename).suffix.lower()
        text_content = ""
        
        if file_ext == '.pdf':
            pdf_reader = PdfReader(io.BytesIO(file_bytes))
            text_content = "\n\n".join([page.extract_text() or "" for page in pdf_reader.pages])
            
        elif file_ext == '.docx':
            doc = Document(io.BytesIO(file_bytes))
            text_content = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
            
        elif file_ext == '.pptx':
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slides_text.append("\n".join(slide_text))
            text_content = "\n\n---SLIDE---\n\n".join(slides_text)
            
        elif file_ext in ['.xlsx', '.xls']:
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            sheets_text = []
            for sheet in wb.worksheets:
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        sheet_data.append(row_text)
                sheets_text.append(f"Sheet: {sheet.title}\n" + "\n".join(sheet_data))
            text_content = "\n\n---SHEET---\n\n".join(sheets_text)
            
        elif file_ext == '.csv':
            try:
                df = pd.read_csv(io.BytesIO(file_bytes))
                text_content = df.to_string(index=False)
            except Exception:
                text_content = file_bytes.decode('utf-8', errors='ignore')
            
        elif file_ext == '.txt':
            text_content = file_bytes.decode('utf-8', errors='ignore')
            
        else:
            text_content = f"[Unsupported file type: {file_ext}]"
        
        if len(text_content) > 50000:
            text_content = text_content[:50000] + "\n\n[... content truncated ...]"
        
        return text_content.strip()
        
    except Exception as e:
        return f"[Error extracting text: {str(e)}]"

# Connect to PostgreSQL
db_conn = psycopg2.connect(
    host="postgres",
    database="docubrain",
    user="docbrain",
    password="docbrain123"
)

# Connect to MinIO
minio_client = Minio(
    "minio:9000",
    access_key="minioadmin",
    secret_key="minioadmin123",
    secure=False
)

# Get all documents without extracted text
with db_conn.cursor() as cur:
    cur.execute("SELECT id, filename, minio_path FROM documents WHERE extracted_text IS NULL")
    documents = cur.fetchall()

print(f"Found {len(documents)} documents to process...")

# Process each document
for doc_id, filename, minio_path in documents:
    try:
        print(f"Processing: {filename}...")
        
        # Download from MinIO
        response = minio_client.get_object("documents", minio_path)
        file_bytes = response.read()
        response.close()
        response.release_conn()
        
        # Extract text
        extracted_text = extract_text_from_file(file_bytes, filename)
        
        # Update database
        with db_conn.cursor() as cur:
            cur.execute(
                "UPDATE documents SET extracted_text = %s WHERE id = %s",
                (extracted_text, doc_id)
            )
        db_conn.commit()
        
        print(f"  ✅ Extracted {len(extracted_text)} characters")
        
    except Exception as e:
        print(f"  ❌ Error: {e}")
        db_conn.rollback()

print("\n✅ Migration complete!")
db_conn.close()
