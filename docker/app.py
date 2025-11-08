try:
    import streamlit as st
except Exception:
    # Minimal stub for environments where streamlit is not installed (prevents import errors
    # and allows static analysis / light-weight execution). This stub returns benign defaults
    # for commonly used Streamlit APIs in this app.
    class _DummyCtx:
        def __init__(self, *args, **kwargs):
            pass
        def __enter__(self):
            return self
        def __exit__(self, exc_type, exc, tb):
            return False
        def __call__(self, *a, **k):
            return self

    class _StStub:
        def __init__(self):
            # used like a dict in the app
            self.session_state = {}
            # used as "with st.sidebar:"
            self.sidebar = _DummyCtx()

        def set_page_config(self, *a, **k):
            return None

        # decorator used as @st.cache_resource
        def cache_resource(self, func=None, **kwargs):
            if func is None:
                def decorator(f):
                    return f
                return decorator
            return func

        # display helpers (no-op)
        def title(self, *a, **k): return None
        def markdown(self, *a, **k): return None
        def subheader(self, *a, **k): return None
        def header(self, *a, **k): return None
        def write(self, *a, **k): return None
        def info(self, *a, **k): return None
        def success(self, *a, **k): return None
        def error(self, *a, **k): return None
        def warning(self, *a, **k): return None
        def stop(self, *a, **k): return None
        def rerun(self, *a, **k): return None
        def divider(self, *a, **k): return None
        def caption(self, *a, **k): return None
        def code(self, *a, **k): return None
        def metric(self, *a, **k): return None

        # input helpers - return safe defaults
        def text_input(self, *a, **k): return ""
        def file_uploader(self, *a, **k): return []
        def button(self, *a, **k): return False

        # layout helpers - return context managers
        def columns(self, specs):
            try:
                length = len(specs) if isinstance(specs, (list, tuple)) else int(specs)
            except Exception:
                length = 1
            return tuple(_DummyCtx() for _ in range(length))

        def tabs(self, names):
            try:
                length = len(names)
            except Exception:
                length = 1
            return tuple(_DummyCtx() for _ in range(length))

    st = _StStub()

import os
from pathlib import Path
from minio import Minio
from minio.error import S3Error
import psycopg2
from datetime import datetime
import io
import hashlib
import threading
import time
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import requests
import json
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
import pandas as pd

# Page config
st.set_page_config(page_title="DocuBrain", page_icon="🧠", layout="wide")

# Session state for user
if 'user_id' not in st.session_state:
    st.session_state.user_id = None
if 'username' not in st.session_state:
    st.session_state.username = None
if 'is_admin' not in st.session_state:
    st.session_state.is_admin = False
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []

# Initialize connections
@st.cache_resource
def init_minio():
    return Minio(
        os.getenv("MINIO_ENDPOINT", "minio:9000"),
        access_key=os.getenv("MINIO_ACCESS_KEY", "minioadmin"),
        secret_key=os.getenv("MINIO_SECRET_KEY", "minioadmin123"),
        secure=False
    )

@st.cache_resource
def init_db():
    conn = psycopg2.connect(
        host=os.getenv("POSTGRES_HOST", "postgres"),
        database=os.getenv("POSTGRES_DB", "docubrain"),
        user=os.getenv("POSTGRES_USER", "docbrain"),
        password=os.getenv("POSTGRES_PASSWORD", "docbrain123")
    )
    conn.autocommit = True

    with conn.cursor() as cur:
        # Create tables if not exists (no destructive drops)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS users (
                id SERIAL PRIMARY KEY,
                username VARCHAR(100) UNIQUE NOT NULL,
                password_hash VARCHAR(255) NOT NULL,
                is_admin BOOLEAN DEFAULT FALSE,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS documents (
                id SERIAL PRIMARY KEY,
                filename VARCHAR(255) NOT NULL,
                file_hash VARCHAR(64) UNIQUE NOT NULL,
                file_size INTEGER NOT NULL,
                file_type VARCHAR(128),
                upload_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                uploaded_by INTEGER REFERENCES users(id),
                status VARCHAR(50) DEFAULT 'uploaded',
                minio_path VARCHAR(255),
                extracted_text TEXT,
                source VARCHAR(50) DEFAULT 'upload'
            )
        """)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS chat_history (
                id SERIAL PRIMARY KEY,
                user_id INTEGER REFERENCES users(id),
                document_id INTEGER REFERENCES documents(id),
                message TEXT NOT NULL,
                response TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        cur.execute("""
            CREATE TABLE IF NOT EXISTS activity_log (
                id SERIAL PRIMARY KEY,
                user_id INTEGER REFERENCES users(id),
                action VARCHAR(100),
                details TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        # Seed users idempotently
        admin_hash = hashlib.sha256("admin".encode()).hexdigest()
        user_hash = hashlib.sha256("user".encode()).hexdigest()
        cur.execute("INSERT INTO users (username, password_hash, is_admin) VALUES ('admin', %s, TRUE) ON CONFLICT (username) DO NOTHING", (admin_hash,))
        cur.execute("INSERT INTO users (username, password_hash, is_admin) VALUES ('user', %s, FALSE) ON CONFLICT (username) DO NOTHING", (user_hash,))

        # Migrate file_type length if still 50
        cur.execute("""
            SELECT character_maximum_length
            FROM information_schema.columns
            WHERE table_name='documents' AND column_name='file_type'
        """)
        row = cur.fetchone()
        if row and (row[0] is None or row[0] < 128):
            cur.execute("ALTER TABLE documents ALTER COLUMN file_type TYPE VARCHAR(128)")
    return conn

minio_client = init_minio()
db_conn = init_db()

# Ensure buckets exist
for bucket in ["documents", "processed"]:
    try:
        if not minio_client.bucket_exists(bucket):
            minio_client.make_bucket(bucket)
    except S3Error:
        pass

# File system watcher
class FolderWatcher(FileSystemEventHandler):
    def on_created(self, event):
        if not event.is_directory:
            time.sleep(1)
            process_watch_file(event.src_path)

@st.cache_resource
def start_folder_watcher():
    watch_folder = Path(os.getenv("WATCH_FOLDER", "/app/watch"))
    watch_folder.mkdir(parents=True, exist_ok=True)
    
    observer = Observer()
    observer.schedule(FolderWatcher(), str(watch_folder), recursive=False)
    observer.start()
    return observer

# Safe helper
def safe_file_type(uploaded_file):
    # Prefer MIME type; fallback to extension; truncate to 120 chars
    mime = (getattr(uploaded_file, "type", None) or "").strip()
    if not mime:
        ext = Path(uploaded_file.name).suffix.lstrip(".").lower() or "unknown"
        mime = ext
    return mime[:120]

# Document text extraction
def extract_text_from_file(file_bytes, filename):
    """Extract text content from various file types"""
    try:
        file_ext = Path(filename).suffix.lower()
        text_content = ""
        
        if file_ext == '.pdf':
            # Extract from PDF
            pdf_reader = PdfReader(io.BytesIO(file_bytes))
            text_content = "\n\n".join([page.extract_text() or "" for page in pdf_reader.pages])
            
        elif file_ext == '.docx':
            # Extract from Word document
            doc = Document(io.BytesIO(file_bytes))
            text_content = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])
            
        elif file_ext == '.pptx':
            # Extract from PowerPoint
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slides_text.append("\n".join(slide_text))
            text_content = "\n\n---SLIDE---\n\n".join(slides_text)
            
        elif file_ext in ['.xlsx', '.xls']:
            # Extract from Excel
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            sheets_text = []
            for sheet in wb.worksheets:
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        sheet_data.append(row_text)
                sheets_text.append(f"Sheet: {sheet.title}\n" + "\n".join(sheet_data))
            text_content = "\n\n---SHEET---\n\n".join(sheets_text)
            
        elif file_ext == '.csv':
            # Extract from CSV
            try:
                df = pd.read_csv(io.BytesIO(file_bytes))
                text_content = df.to_string(index=False)
            except Exception:
                text_content = file_bytes.decode('utf-8', errors='ignore')
            
        elif file_ext == '.txt':
            # Extract from text file
            text_content = file_bytes.decode('utf-8', errors='ignore')
            
        else:
            # Unsupported file type
            text_content = f"[Unsupported file type: {file_ext}]"
        
        # Truncate if too long (keep first 50k characters to avoid DB issues)
        if len(text_content) > 50000:
            text_content = text_content[:50000] + "\n\n[... content truncated ...]"
        
        return text_content.strip()
        
    except Exception as e:
        return f"[Error extracting text: {str(e)}]"

def process_watch_file(filepath):
    """Process files dropped in watch folder"""
    try:
        filepath = Path(filepath)
        if filepath.suffix.lower() not in ['.pdf', '.txt', '.docx', '.pptx', '.xlsx', '.csv']:
            return
        
        with open(filepath, 'rb') as f:
            file_bytes = f.read()
        
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        file_size = len(file_bytes)
        
        with db_conn.cursor() as cur:
            cur.execute("SELECT id FROM documents WHERE file_hash = %s", (file_hash,))
            if cur.fetchone():
                filepath.unlink()
                return
        
        minio_path = f"watch/{datetime.now().strftime('%Y/%m/%d')}/{file_hash[:8]}_{filepath.name}"
        minio_client.put_object(
            "documents",
            minio_path,
            io.BytesIO(file_bytes),
            file_size
        )
        
        # Extract text content from the file
        extracted_text = extract_text_from_file(file_bytes, filepath.name)
        
        try:
            with db_conn.cursor() as cur:
                cur.execute("""
                    INSERT INTO documents (filename, file_hash, file_size, file_type, minio_path, source, extracted_text)
                    VALUES (%s, %s, %s, %s, %s, 'folder_watch', %s)
                """, (filepath.name, file_hash, file_size, (filepath.suffix or '').lstrip('.')[:120], minio_path, extracted_text))
        except Exception as e:
            try:
                db_conn.rollback()
            except:
                pass
            print(f"Error processing {filepath}: {e}")
            return
        
        filepath.unlink()
        
    except Exception as e:
        print(f"Error processing {filepath}: {e}")

# Ollama AI Integration
def query_ollama(prompt, context="", model="llama3:latest"):
    """Query Ollama API for AI responses"""
    try:
        # Try to connect to Ollama on host machine (Windows)
        ollama_url = os.getenv("OLLAMA_URL", "http://host.docker.internal:11434/api/generate")
        
        # Shorter prompt for faster models
        if model in ["tinyllama", "phi3:mini", "llama3.2:3b"]:
            full_prompt = f"""Context: {context[:800] if context else 'No documents'}

Question: {prompt}

Answer briefly:"""
            max_tokens = 100  # Very short for fast response
        else:
            full_prompt = f"""You are DocuBrain, an AI assistant helping users understand their documents.
        
Context from documents:
{context}

User question: {prompt}

Please provide a helpful, concise answer based on the context provided. If the context doesn't contain relevant information, say so politely."""
            max_tokens = 150

        payload = {
            "model": model,
            "prompt": full_prompt,
            "stream": False,
            "options": {
                "temperature": 0.7,
                "num_predict": max_tokens,
                "top_k": 40,
                "top_p": 0.9,
                "num_ctx": 2048  # Reduced context window for speed
            }
        }
        
        # Timeout: 300s for fast models, 600s for larger ones
        timeout = 300 if model in ["tinyllama", "phi3:mini", "llama3.2:3b"] else 600
        response = requests.post(ollama_url, json=payload, timeout=timeout)
        
        if response.status_code == 200:
            result = response.json().get("response", "I apologize, but I couldn't generate a response.")
            return result.strip()
        elif response.status_code == 404:
            return f"⚠️ Model '{model}' not found. Please pull it with: ollama pull {model}"
        else:
            return f"⚠️ Ollama error (Status {response.status_code}). Please check if Ollama is running."
    
    except requests.exceptions.ConnectionError:
        return "⚠️ Cannot connect to Ollama. Please ensure Ollama is running on your system (http://localhost:11434)"
    except requests.exceptions.Timeout:
        return "⚠️ Request timed out. The model might be loading or the query is too complex."
    except Exception as e:
        return f"⚠️ Error: {str(e)}"

def get_document_context(user_id, query, limit=3):
    """Get relevant document context for AI query"""
    try:
        with db_conn.cursor() as cur:
            # Get recent documents from the user or all documents if admin uploaded
            cur.execute("""
                SELECT filename, extracted_text, file_type, upload_date
                FROM documents
                WHERE (uploaded_by = %s OR uploaded_by IS NULL) 
                  AND extracted_text IS NOT NULL
                  AND extracted_text != '[Unsupported file type: .pdf]'
                  AND extracted_text != '[Error extracting text: '
                ORDER BY upload_date DESC
                LIMIT %s
            """, (user_id, limit))
            
            docs = cur.fetchall()
            
            if not docs:
                return "No documents with extracted text available yet. Please upload some documents (PDF, DOCX, TXT, etc.)."
            
            # Build context with more information
            context_parts = []
            for filename, text, file_type, upload_date in docs:
                # Reduced to 800 chars for faster processing
                text_excerpt = (text or '').strip()
                if len(text_excerpt) > 800:
                    text_excerpt = text_excerpt[:800] + "..."
                
                context_parts.append(
                    f"📄 {filename} ({file_type})\n{text_excerpt}\n"
                )
            
            return "\n---\n".join(context_parts)
    except Exception as e:
        return f"Error retrieving documents: {str(e)}"

# Start watcher
watcher = start_folder_watcher()

# Login/User Management
def authenticate_user(username, password):
    """Authenticate user with username and password"""
    password_hash = hashlib.sha256(password.encode()).hexdigest()
    
    with db_conn.cursor() as cur:
        cur.execute("""
            SELECT id, username, is_admin 
            FROM users 
            WHERE username = %s AND password_hash = %s
        """, (username, password_hash))
        result = cur.fetchone()
    
    if result:
        user_id, username, is_admin = result
        
        # Log activity
        with db_conn.cursor() as cur:
            cur.execute("INSERT INTO activity_log (user_id, action) VALUES (%s, %s)", 
                       (user_id, "login"))
            db_conn.commit()
        
        return user_id, username, is_admin
    
    return None, None, False

def register_user(username, password):
    """Register a new user"""
    password_hash = hashlib.sha256(password.encode()).hexdigest()
    
    try:
        with db_conn.cursor() as cur:
            cur.execute("""
                INSERT INTO users (username, password_hash, is_admin)
                VALUES (%s, %s, FALSE)
                RETURNING id
            """, (username, password_hash))
            user_id = cur.fetchone()[0]
            db_conn.commit()
            
            # Log activity
            cur.execute("INSERT INTO activity_log (user_id, action) VALUES (%s, %s)", 
                       (user_id, "register"))
            db_conn.commit()
        
        return user_id, username, False
    except psycopg2.IntegrityError:
        db_conn.rollback()
        return None, None, False

# Login UI
if st.session_state.user_id is None:
    st.title("🧠 DocuBrain")
    st.markdown("### Document Intelligence System")
    
    col1, col2, col3 = st.columns([1, 2, 1])
    
    with col2:
        st.markdown("---")
        st.subheader("🔐 Login")
        
        tab1, tab2 = st.tabs(["Login", "Register"])
        
        with tab1:
            username = st.text_input("Username", key="login_username")
            password = st.text_input("Password", type="password", key="login_password")
            
            col_a, col_b = st.columns(2)
            with col_a:
                if st.button("Login", use_container_width=True):
                    if username and password:
                        user_id, user_name, is_admin = authenticate_user(username, password)
                        if user_id:
                            st.session_state.user_id = user_id
                            st.session_state.username = user_name
                            st.session_state.is_admin = is_admin
                            st.success(f"Welcome, {user_name}!")
                            st.rerun()
                        else:
                            st.error("Invalid username or password")
                    else:
                        st.warning("Please enter username and password")
            
            with col_b:
                st.info("**Default Admin:**\nuser: `admin`\npass: `admin`")
        
        with tab2:
            new_username = st.text_input("Username", key="reg_username")
            new_password = st.text_input("Password", type="password", key="reg_password")
            confirm_password = st.text_input("Confirm Password", type="password", key="reg_confirm")
            
            if st.button("Register", use_container_width=True):
                if new_username and new_password and confirm_password:
                    if new_password == confirm_password:
                        if len(new_password) >= 4:
                            user_id, user_name, is_admin = register_user(new_username, new_password)
                            if user_id:
                                st.session_state.user_id = user_id
                                st.session_state.username = user_name
                                st.session_state.is_admin = is_admin
                                st.success(f"Registration successful! Welcome, {user_name}!")
                                st.rerun()
                            else:
                                st.error("Username already exists")
                        else:
                            st.warning("Password must be at least 4 characters")
                    else:
                        st.error("Passwords don't match")
                else:
                    st.warning("Please fill all fields")
    
    st.stop()

# Sidebar - User Info
with st.sidebar:
    st.header("👤 User")
    
    st.success(f"**{st.session_state.username}**")
    if st.session_state.is_admin:
        st.info("🔑 Administrator")
    
    if st.button("Logout", use_container_width=True):
        st.session_state.user_id = None
        st.session_state.username = None
        st.session_state.is_admin = False
        st.session_state.chat_history = []
        st.rerun()
    
    st.divider()
    
    # System Status
    st.header("System Status")
    
    try:
        with db_conn.cursor() as cur:
            cur.execute("SELECT 1")
        st.success("✅ PostgreSQL")
    except:
        st.error("❌ PostgreSQL")
    
    try:
        minio_client.bucket_exists("documents")
        st.success("✅ MinIO")
    except:
        st.error("❌ MinIO")
    
    st.divider()
    
    # Statistics
    st.header("📊 Statistics")
    try:
        with db_conn.cursor() as cur:
            cur.execute("SELECT COUNT(*) FROM documents")
            doc_count = cur.fetchone()[0]
            st.metric("Total Documents", doc_count)
    except Exception as e:
        db_conn.rollback()
        st.metric("Total Documents", "Error")
    
    try:
        with db_conn.cursor() as cur:
            cur.execute("SELECT COUNT(DISTINCT uploaded_by) FROM documents WHERE uploaded_by IS NOT NULL")
            user_count = cur.fetchone()[0]
            st.metric("Active Users", user_count)
    except Exception as e:
        db_conn.rollback()
        st.metric("Active Users", "Error")
    
    try:
        with db_conn.cursor() as cur:
            cur.execute("SELECT COUNT(*) FROM documents WHERE source = 'folder_watch'")
            watch_count = cur.fetchone()[0]
            st.metric("Auto-imported", watch_count)
    except Exception as e:
        db_conn.rollback()
        st.metric("Auto-imported", "Error")

    st.divider()
    
    # Network access info
    st.header("🌐 Network Access")
    try:
        hostname = os.popen('hostname -I').read().strip().split()[0]
    except:
        hostname = 'YOUR_IP'
    st.code(f"http://{hostname}:8501")
    st.caption("Share this URL with your team")

# Main UI
st.title("🧠 DocuBrain")
st.markdown(f"### Welcome, {st.session_state.username}!")

# Add custom CSS for better tab styling
st.markdown("""
<style>
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
    }
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        padding-left: 20px;
        padding-right: 20px;
        background-color: rgba(28, 131, 225, 0.1);
        border-radius: 5px;
    }
    .stTabs [aria-selected="true"] {
        background-color: rgba(28, 131, 225, 0.3);
    }
</style>
""", unsafe_allow_html=True)

tab1, tab2, tab3, tab4, tab5 = st.tabs([
    "📤 Upload Documents", 
    "📚 Document Library", 
    "🔍 Search & Filter", 
    "💬 AI Chat", 
    "📊 Activity Log"
])

with tab1:
    st.header("Upload Documents")
    
    # Upload method selection
    upload_method = st.radio(
        "Upload Method",
        ["📄 Individual Files", "📁 Bulk Folder Import", "☁️ Import from Storage"],
        horizontal=True
    )
    
    if upload_method == "📄 Individual Files":
        st.subheader("Upload Individual Files")
        uploaded_files = st.file_uploader(
            "Choose files",
            type=["pdf", "txt", "docx", "pptx", "xlsx", "csv", "md"],
            accept_multiple_files=True
        )
        
        if uploaded_files:
            for uploaded_file in uploaded_files:
                file_bytes = uploaded_file.read()
                file_hash = hashlib.sha256(file_bytes).hexdigest()
                file_size = len(file_bytes)
                
                with db_conn.cursor() as cur:
                    cur.execute("SELECT filename FROM documents WHERE file_hash = %s", (file_hash,))
                    existing = cur.fetchone()
                
                if existing:
                    st.warning(f"⚠️ {uploaded_file.name} already exists")
                    continue
                
                try:
                    minio_path = f"uploads/{datetime.now().strftime('%Y/%m/%d')}/{file_hash[:8]}_{uploaded_file.name}"
                    minio_client.put_object(
                        "documents",
                        minio_path,
                        io.BytesIO(file_bytes),
                        file_size
                    )

                    # Extract text content from the file
                    extracted_text = extract_text_from_file(file_bytes, uploaded_file.name)

                    ft = safe_file_type(uploaded_file)
                    with db_conn.cursor() as cur:
                        cur.execute("""
                            INSERT INTO documents (filename, file_hash, file_size, file_type, minio_path, uploaded_by, extracted_text)
                            VALUES (%s, %s, %s, %s, %s, %s, %s)
                        """, (uploaded_file.name, file_hash, file_size, ft, minio_path, st.session_state.user_id, extracted_text))
                        # autocommit=True, commit not required

                        cur.execute("INSERT INTO activity_log (user_id, action, details) VALUES (%s, %s, %s)",
                                    (st.session_state.user_id, "upload", uploaded_file.name))
                    st.success(f"✅ {uploaded_file.name} ({file_size / 1024:.1f} KB) - Text extracted!")
                except Exception as e:
                    # Clear any transaction state and report
                    try:
                        db_conn.rollback()
                    except:
                        pass
                    st.error(f"❌ Failed: {e}")
    
    elif upload_method == "📁 Bulk Folder Import":
        st.subheader("📁 Bulk Folder Import - Drop Zone Method")
        
        # Get watch folder path
        watch_folder = os.path.join(os.path.expanduser("~"), "DocuBrain", "watch")
        
        st.info(f"""**Drop Zone Location:**
        
📁 `{watch_folder}`
        
**How to upload an entire folder:**
1. Open File Explorer
2. Find your folder (e.g., C:\\\\MyDocuments\\\\ProjectFiles)
3. **Copy or drag the ENTIRE folder** into the drop zone above
4. Click "🔄 Import All Files" button below
5. All files (including subfolders) will be imported!
        """)
        
        st.success("""**✅ Advantages:**
- Upload ENTIRE folders (not just individual files)
- Includes all subfolders recursively  
- No file selection needed
- No browser limitations
- Works with hundreds or thousands of files
        """)
        
        # Create watch folder if it doesn't exist
        os.makedirs(watch_folder, exist_ok=True)
        
        # Scan the watch folder for files
        if st.button("🔄 Scan & Import All Files from Drop Zone", type="primary", use_container_width=True):
            import glob
            
            # Find all supported files recursively
            supported_extensions = ['*.pdf', '*.txt', '*.docx', '*.pptx', '*.xlsx', '*.csv', '*.md']
            found_files = []
            
            for ext in supported_extensions:
                found_files.extend(glob.glob(os.path.join(watch_folder, '**', ext), recursive=True))
                found_files.extend(glob.glob(os.path.join(watch_folder, '**', ext.upper()), recursive=True))
            
            if not found_files:
                st.warning(f"⚠️ No files found in the drop zone: {watch_folder}")
                st.info("� **Tip:** Copy your entire folder (with all its contents) to the drop zone path shown above, then click this button again.")
            else:
                st.info(f"📊 Found **{len(found_files)} files** in drop zone. Starting import...")
                
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                success_count = 0
                skip_count = 0
                error_count = 0
                
                for idx, file_path in enumerate(found_files):
                    filename = os.path.basename(file_path)
                    status_text.text(f"Processing {idx + 1}/{len(found_files)}: {filename}")
                    
                    try:
                        # Read file
                        with open(file_path, 'rb') as f:
                            file_bytes = f.read()
                        
                        file_hash = hashlib.sha256(file_bytes).hexdigest()
                        file_size = len(file_bytes)
                        
                        # Check for duplicates
                        with db_conn.cursor() as cur:
                            cur.execute("SELECT filename FROM documents WHERE file_hash = %s", (file_hash,))
                            existing = cur.fetchone()
                        
                        if existing:
                            skip_count += 1
                            # Delete the file since it's a duplicate
                            try:
                                os.remove(file_path)
                            except:
                                pass
                            continue
                        
                        # Upload to MinIO
                        minio_path = f"uploads/{datetime.now().strftime('%Y/%m/%d')}/{file_hash[:8]}_{filename}"
                        minio_client.put_object(
                            "documents",
                            minio_path,
                            io.BytesIO(file_bytes),
                            file_size
                        )
                        
                        # Extract text
                        extracted_text = extract_text_from_file(file_bytes, filename)
                        
                        # Save to database
                        file_ext = filename.split('.')[-1].lower() if '.' in filename else 'unknown'
                        with db_conn.cursor() as cur:
                            cur.execute("""
                                INSERT INTO documents (filename, file_hash, file_size, file_type, minio_path, uploaded_by, extracted_text, source)
                                VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                            """, (filename, file_hash, file_size, file_ext, minio_path, st.session_state.user_id, extracted_text, "bulk_import"))
                            
                            cur.execute("INSERT INTO activity_log (user_id, action, details) VALUES (%s, %s, %s)",
                                        (st.session_state.user_id, "bulk_import", filename))
                        
                        # Delete the file after successful upload
                        try:
                            os.remove(file_path)
                        except:
                            pass
                        
                        success_count += 1
                    
                    except Exception as e:
                        error_count += 1
                        try:
                            db_conn.rollback()
                        except:
                            pass
                    
                    progress_bar.progress((idx + 1) / len(found_files))
                
                status_text.empty()
                progress_bar.empty()
                
                st.success(f"""
### 📊 Import Complete!
- ✅ **Imported**: {success_count} files
- ⏭️ **Skipped** (duplicates): {skip_count} files
- ❌ **Errors**: {error_count} files
- 📦 **Total processed**: {len(found_files)} files
- 📁 **Source**: Drop Zone ({watch_folder})
                """)
                
                if success_count > 0:
                    st.balloons()
                
                # Clean up empty folders
                try:
                    for root, dirs, files in os.walk(watch_folder, topdown=False):
                        for dir_name in dirs:
                            dir_path = os.path.join(root, dir_name)
                            try:
                                os.rmdir(dir_path)  # Only removes if empty
                            except:
                                pass
                except:
                    pass
        
        # Show current files in watch folder
        st.divider()
        with st.expander("📋 View Files Currently in Drop Zone"):
            import glob
            supported_extensions = ['*.pdf', '*.txt', '*.docx', '*.pptx', '*.xlsx', '*.csv', '*.md']
            current_files = []
            
            for ext in supported_extensions:
                current_files.extend(glob.glob(os.path.join(watch_folder, '**', ext), recursive=True))
                current_files.extend(glob.glob(os.path.join(watch_folder, '**', ext.upper()), recursive=True))
            
            if current_files:
                st.write(f"**Found {len(current_files)} files:**")
                for file_path in current_files[:20]:
                    rel_path = os.path.relpath(file_path, watch_folder)
                    st.text(f"📄 {rel_path}")
                if len(current_files) > 20:
                    st.text(f"... and {len(current_files) - 20} more files")
            else:
                st.info("No files currently in drop zone. Copy your folder here to get started!")

            
            # Process scanned files
            if 'scanned_files' in st.session_state and st.session_state.scanned_files:
                st.write(f"### 📊 Found {len(st.session_state.scanned_files)} files")
                
                # Show preview of first 10 files
                with st.expander("📋 Preview files"):
                    for file_path in st.session_state.scanned_files[:10]:
                        st.text(f"📄 {os.path.basename(file_path)}")
                    if len(st.session_state.scanned_files) > 10:
                        st.text(f"... and {len(st.session_state.scanned_files) - 10} more files")
                
                organize_by_date = st.checkbox("📅 Organize by upload date", value=True, key="local_folder_organize")
                
                if st.button("🚀 Upload All Files from Folder", type="primary", key="upload_local_folder"):
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    success_count = 0
                    skip_count = 0
                    error_count = 0
                    
                    for idx, file_path in enumerate(st.session_state.scanned_files):
                        filename = os.path.basename(file_path)
                        status_text.text(f"Processing {idx + 1}/{len(st.session_state.scanned_files)}: {filename}")
                        
                        try:
                            # Read file from disk
                            with open(file_path, 'rb') as f:
                                file_bytes = f.read()
                            
                            file_hash = hashlib.sha256(file_bytes).hexdigest()
                            file_size = len(file_bytes)
                            
                            # Check for duplicates
                            with db_conn.cursor() as cur:
                                cur.execute("SELECT filename FROM documents WHERE file_hash = %s", (file_hash,))
                                existing = cur.fetchone()
                            
                            if existing:
                                skip_count += 1
                                continue
                            
                            # Organize path
                            if organize_by_date:
                                minio_path = f"uploads/{datetime.now().strftime('%Y/%m/%d')}/{file_hash[:8]}_{filename}"
                            else:
                                minio_path = f"uploads/{file_hash[:8]}_{filename}"
                            
                            # Upload to MinIO
                            minio_client.put_object(
                                "documents",
                                minio_path,
                                io.BytesIO(file_bytes),
                                file_size
                            )
                            
                            # Extract text
                            extracted_text = extract_text_from_file(file_bytes, filename)
                            
                            # Save to database
                            file_ext = filename.split('.')[-1].lower() if '.' in filename else 'unknown'
                            with db_conn.cursor() as cur:
                                cur.execute("""
                                    INSERT INTO documents (filename, file_hash, file_size, file_type, minio_path, uploaded_by, extracted_text, source)
                                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                                """, (filename, file_hash, file_size, file_ext, minio_path, st.session_state.user_id, extracted_text, "folder_path"))
                                
                                cur.execute("INSERT INTO activity_log (user_id, action, details) VALUES (%s, %s, %s)",
                                            (st.session_state.user_id, "folder_upload", filename))
                            
                            success_count += 1
                        
                        except Exception as e:
                            error_count += 1
                            st.warning(f"⚠️ {filename}: {str(e)}")
                            try:
                                db_conn.rollback()
                            except:
                                pass
                        
                        progress_bar.progress((idx + 1) / len(st.session_state.scanned_files))
                    
                    status_text.empty()
                    progress_bar.empty()
                    
                    st.success(f"""
                    ### 📊 Upload Complete!
                    - ✅ **Uploaded**: {success_count} files
                    - ⏭️ **Skipped** (duplicates): {skip_count} files
                    - ❌ **Errors**: {error_count} files
                    - 📦 **Total processed**: {len(st.session_state.scanned_files)} files
                    - 📁 **Source folder**: {st.session_state.scan_folder_path}
                    """)
                    
                    if success_count > 0:
                        st.balloons()
                    
                    # Clear session state
                    del st.session_state.scanned_files
                    del st.session_state.scan_folder_path
            
            folder_files = None  # Set to None so the web upload section doesn't run
        
        if folder_files:
            st.write(f"📊 **Selected {len(folder_files)} files**")
            
            # Option to organize by subfolder
            organize_by_date = st.checkbox("📅 Organize by upload date", value=True)
            
            # Process button
            if st.button("🚀 Upload All Files", type="primary"):
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                success_count = 0
                skip_count = 0
                error_count = 0
                
                for idx, uploaded_file in enumerate(folder_files):
                    status_text.text(f"Processing {idx + 1}/{len(folder_files)}: {uploaded_file.name}")
                    
                    try:
                        file_bytes = uploaded_file.read()
                        file_hash = hashlib.sha256(file_bytes).hexdigest()
                        file_size = len(file_bytes)
                        
                        # Check for duplicates
                        with db_conn.cursor() as cur:
                            cur.execute("SELECT filename FROM documents WHERE file_hash = %s", (file_hash,))
                            existing = cur.fetchone()
                        
                        if existing:
                            skip_count += 1
                            continue
                        
                        # Organize path
                        if organize_by_date:
                            minio_path = f"uploads/{datetime.now().strftime('%Y/%m/%d')}/{file_hash[:8]}_{uploaded_file.name}"
                        else:
                            minio_path = f"uploads/{file_hash[:8]}_{uploaded_file.name}"
                        
                        # Upload to MinIO
                        minio_client.put_object(
                            "documents",
                            minio_path,
                            io.BytesIO(file_bytes),
                            file_size
                        )
                        
                        # Extract text
                        extracted_text = extract_text_from_file(file_bytes, uploaded_file.name)
                        
                        # Save to database
                        ft = safe_file_type(uploaded_file)
                        with db_conn.cursor() as cur:
                            cur.execute("""
                                INSERT INTO documents (filename, file_hash, file_size, file_type, minio_path, uploaded_by, extracted_text, source)
                                VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                            """, (uploaded_file.name, file_hash, file_size, ft, minio_path, st.session_state.user_id, extracted_text, "folder_upload"))
                            
                            cur.execute("INSERT INTO activity_log (user_id, action, details) VALUES (%s, %s, %s)",
                                        (st.session_state.user_id, "folder_upload", uploaded_file.name))
                        
                        success_count += 1
                    
                    except Exception as e:
                        error_count += 1
                        st.warning(f"⚠️ {uploaded_file.name}: {str(e)}")
                        try:
                            db_conn.rollback()
                        except:
                            pass
                    
                    # Update progress
                    progress_bar.progress((idx + 1) / len(folder_files))
                
                # Summary
                status_text.empty()
                progress_bar.empty()
                
                st.success(f"""
                ### 📊 Upload Complete!
                - ✅ **Uploaded**: {success_count} files
                - ⏭️ **Skipped** (duplicates): {skip_count} files
                - ❌ **Errors**: {error_count} files
                - 📦 **Total processed**: {len(folder_files)} files
                """)
                
                if success_count > 0:
                    st.balloons()
    
    elif upload_method == "☁️ Import from Storage":
        st.subheader("Import from MinIO Storage")
        st.info("📦 Browse and import files already stored in MinIO but not yet registered in the database.")
        
        try:
            # List all objects in the documents bucket
            objects = minio_client.list_objects("documents", recursive=True)
            
            # Get all file hashes already in database
            with db_conn.cursor() as cur:
                cur.execute("SELECT file_hash FROM documents")
                existing_hashes = {row[0] for row in cur.fetchall()}
            
            # Filter for unregistered files
            unregistered_files = []
            all_files = []
            
            for obj in objects:
                all_files.append({
                    'name': obj.object_name,
                    'size': obj.size,
                    'last_modified': obj.last_modified
                })
                
                # Try to extract hash from filename (format: hash[:8]_filename)
                filename_parts = obj.object_name.split('/')[-1]
                if '_' in filename_parts:
                    potential_hash = filename_parts.split('_')[0]
                    # Check if this file is in database
                    with db_conn.cursor() as cur:
                        cur.execute("SELECT id FROM documents WHERE minio_path = %s", (obj.object_name,))
                        if not cur.fetchone():
                            unregistered_files.append({
                                'path': obj.object_name,
                                'name': filename_parts,
                                'size': obj.size,
                                'last_modified': obj.last_modified
                            })
            
            st.metric("Total files in storage", len(all_files))
            st.metric("Unregistered files", len(unregistered_files))
            
            if unregistered_files:
                st.warning(f"⚠️ Found {len(unregistered_files)} files in storage that are not registered in database")
                
                # Show unregistered files
                with st.expander("📋 View Unregistered Files"):
                    for file_info in unregistered_files[:20]:  # Show first 20
                        st.text(f"📄 {file_info['name']} ({file_info['size'] / 1024:.1f} KB)")
                
                if st.button("🔄 Import All Unregistered Files", type="primary"):
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    success_count = 0
                    error_count = 0
                    
                    for idx, file_info in enumerate(unregistered_files):
                        status_text.text(f"Importing {idx + 1}/{len(unregistered_files)}: {file_info['name']}")
                        
                        try:
                            # Download file from MinIO
                            response = minio_client.get_object("documents", file_info['path'])
                            file_bytes = response.read()
                            response.close()
                            
                            file_hash = hashlib.sha256(file_bytes).hexdigest()
                            
                            # Extract text
                            extracted_text = extract_text_from_file(file_bytes, file_info['name'])
                            
                            # Determine file type
                            file_ext = file_info['name'].split('.')[-1].lower() if '.' in file_info['name'] else 'unknown'
                            
                            # Insert into database
                            with db_conn.cursor() as cur:
                                cur.execute("""
                                    INSERT INTO documents (filename, file_hash, file_size, file_type, minio_path, uploaded_by, extracted_text, source)
                                    VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                                """, (file_info['name'], file_hash, file_info['size'], file_ext, file_info['path'], st.session_state.user_id, extracted_text, "minio_import"))
                                
                                cur.execute("INSERT INTO activity_log (user_id, action, details) VALUES (%s, %s, %s)",
                                            (st.session_state.user_id, "minio_import", file_info['name']))
                            
                            success_count += 1
                        
                        except Exception as e:
                            error_count += 1
                            st.warning(f"⚠️ {file_info['name']}: {str(e)}")
                            try:
                                db_conn.rollback()
                            except:
                                pass
                        
                        progress_bar.progress((idx + 1) / len(unregistered_files))
                    
                    status_text.empty()
                    progress_bar.empty()
                    
                    st.success(f"""
                    ### 📊 Import Complete!
                    - ✅ **Imported**: {success_count} files
                    - ❌ **Errors**: {error_count} files
                    - 📦 **Total processed**: {len(unregistered_files)} files
                    """)
                    
                    if success_count > 0:
                        st.balloons()
            else:
                st.success("✅ All files in storage are already registered in the database!")
        
        except Exception as e:
            st.error(f"❌ Error accessing storage: {str(e)}")
    
    else:  # Folder Watch Info
        st.subheader("Folder Watch (Auto-Import)")
        st.info(f"""
        **Auto-import files from:**
        
        📁 `C:\\Users\\{os.getenv('USERNAME', 'USER')}\\DocuBrain\\watch`
        
        **Supported formats:**
        PDF, DOCX, PPTX, XLSX, TXT, CSV
        
        **How it works:**
        1. Drop files into the watch folder
        2. System auto-processes them
        3. Files appear in Library
        4. Original files are deleted after import
        """)
        
        if st.button("🔄 Refresh Now"):
            st.rerun()

with tab2:
    st.header("Document Library")
    
    # Filters
    col1, col2, col3 = st.columns(3)
    with col1:
        filter_source = st.selectbox("Source", ["All", "Upload", "Folder Watch"])
    with col2:
        try:
            with db_conn.cursor() as cur:
                cur.execute("SELECT DISTINCT u.username FROM users u INNER JOIN documents d ON u.id = d.uploaded_by ORDER BY u.username")
                user_list = ["All"] + [row[0] for row in cur.fetchall()]
        except Exception as e:
            db_conn.rollback()
            user_list = ["All"]
        
        filter_user = st.selectbox("User", user_list)
    with col3:
        sort_by = st.selectbox("Sort by", ["Recent", "Name", "Size"])
    
    # Fetch documents
    query = "SELECT d.id, d.filename, d.file_size, d.upload_date, d.source, u.username FROM documents d LEFT JOIN users u ON d.uploaded_by = u.id WHERE 1=1"
    params = []
    
    if filter_source != "All":
        query += " AND d.source = %s"
        params.append("upload" if filter_source == "Upload" else "folder_watch")
    
    if filter_user != "All":
        query += " AND u.username = %s"
        params.append(filter_user)
    
    if sort_by == "Recent":
        query += " ORDER BY d.upload_date DESC"
    elif sort_by == "Name":
        query += " ORDER BY d.filename"
    else:
        query += " ORDER BY d.file_size DESC"
    
    with db_conn.cursor() as cur:
        cur.execute(query, params)
        documents = cur.fetchall()
    
    if documents:
        for doc_id, filename, file_size, upload_date, source, uploader in documents:
            col1, col2, col3, col4, col5 = st.columns([3, 1, 1, 1, 1])
            
            with col1:
                st.write(f"**{filename}**")
            with col2:
                st.write(f"{file_size / 1024:.1f} KB")
            with col3:
                st.write(upload_date.strftime("%Y-%m-%d"))
            with col4:
                st.write(f"{'📤' if source == 'upload' else '📁'} {uploader or 'Auto'}")
            with col5:
                # Only allow deletion if user is admin or owner
                if st.session_state.is_admin or uploader == st.session_state.username:
                    if st.button("🗑️", key=f"del_{doc_id}"):
                        with db_conn.cursor() as cur:
                            cur.execute("DELETE FROM documents WHERE id = %s", (doc_id,))
                            db_conn.commit()
                        st.rerun()
    else:
        st.info("No documents yet.")

with tab3:
    st.header("Search Documents")
    
    search_query = st.text_input("Search filename...")
    
    if search_query:
        with db_conn.cursor() as cur:
            cur.execute("""
                SELECT d.filename, d.file_size, d.upload_date, u.username
                FROM documents d
                LEFT JOIN users u ON d.uploaded_by = u.id
                WHERE d.filename ILIKE %s
                ORDER BY d.upload_date DESC
            """, (f"%{search_query}%",))
            results = cur.fetchall()
        
        if results:
            for filename, file_size, upload_date, uploader in results:
                st.write(f"📄 **{filename}** - {file_size / 1024:.1f} KB - {upload_date.strftime('%Y-%m-%d')} - by {uploader or 'Auto'}")
        else:
            st.info("No results")

with tab4:
    st.header("💬 Chat with Documents")
    
    # AI Model Selection
    col1, col2 = st.columns([3, 1])
    with col1:
        st.info("⚡ Fast models (30-60s): TinyLlama, Phi-3, Llama3.2-3B | Slower models (3-5min): Others")
    with col2:
        ai_model = st.selectbox("Model", [
            "tinyllama",           # ⚡ FASTEST - 1.1B params - Basic but very fast
            "phi3:mini",           # ⚡ FAST - 3.8B params - Great balance
            "llama3.2:3b",         # ⚡ FAST - 3B params - Strong performance
            "qwen2.5:7b",          # Medium - 7B params - Good quality
            "llama3:latest",       # Slow - 8B params - High quality
            "deepseek-r1:8b",      # Slow - 8B params - Reasoning focused
            "granite3.1-dense:8b", # Slow - 8B params - Enterprise model
            "aya-expanse:latest",  # Slow - Multilingual
            "gemma3:12b"           # SLOWEST - 12B params - Best quality
        ], label_visibility="collapsed")
    
    with db_conn.cursor() as cur:
        cur.execute("""
            SELECT message, response, created_at 
            FROM chat_history 
            WHERE user_id = %s 
            ORDER BY created_at DESC 
            LIMIT 10
        """, (st.session_state.user_id,))
        history = cur.fetchall()
    
    if history:
        st.subheader("Recent Conversations")
        for msg, resp, created_at in reversed(history):
            st.markdown(f"**You ({created_at.strftime('%H:%M')}):** {msg}")
            if resp:
                st.markdown(f"**🧠 DocuBrain:** {resp}")
            st.divider()
    
    user_message = st.text_input("Ask a question about your documents...")
    
    col_send, col_info = st.columns([1, 4])
    with col_send:
        send_button = st.button("Send", type="primary", use_container_width=True)
    with col_info:
        # Dynamic timing based on model
        if ai_model in ["tinyllama", "phi3:mini", "llama3.2:3b"]:
            st.caption("⚡ Fast model: 10-30 seconds expected")
        else:
            st.caption("⏳ Large model: 1-5 minutes expected")
    
    if send_button and user_message:
        # Dynamic timing message based on model
        if ai_model in ["tinyllama", "phi3:mini", "llama3.2:3b"]:
            spinner_msg = f"🤔 Generating response with {ai_model}... Should take 10-30 seconds..."
        else:
            spinner_msg = f"🤔 Generating response with {ai_model}... This may take 1-5 minutes..."
        
        with st.spinner(spinner_msg):
            import time
            start_time = time.time()
            
            # Get document context
            context = get_document_context(st.session_state.user_id, user_message)
            
            # Query Ollama
            ai_response = query_ollama(user_message, context, model=ai_model)
            
            elapsed_time = time.time() - start_time
            
            # Save to database
            with db_conn.cursor() as cur:
                cur.execute("""
                    INSERT INTO chat_history (user_id, message, response)
                    VALUES (%s, %s, %s)
                """, (st.session_state.user_id, user_message, ai_response))
                db_conn.commit()
        
        st.success(f"✅ Response generated in {elapsed_time:.1f} seconds!")
        st.rerun()

with tab5:
    st.header("📊 Activity Log")
    
    # Only admins can see all activity
    if st.session_state.is_admin:
        try:
            with db_conn.cursor() as cur:
                cur.execute("""
                    SELECT u.username, a.action, a.details, a.created_at
                    FROM activity_log a
                    JOIN users u ON a.user_id = u.id
                    ORDER BY a.created_at DESC
                    LIMIT 50
                """)
                activities = cur.fetchall()
            
            for username, action, details, created_at in activities:
                st.write(f"`{created_at.strftime('%Y-%m-%d %H:%M')}` **{username}** - {action} {details or ''}")
        except Exception as e:
            db_conn.rollback()
            st.error(f"Error loading activity log: {e}")
    else:
        # Regular users only see their own activity
        try:
            with db_conn.cursor() as cur:
                cur.execute("""
                    SELECT action, details, created_at
                    FROM activity_log
                    WHERE user_id = %s
                    ORDER BY created_at DESC
                    LIMIT 20
                """, (st.session_state.user_id,))
                activities = cur.fetchall()
            
            for action, details, created_at in activities:
                st.write(f"`{created_at.strftime('%Y-%m-%d %H:%M')}` {action} {details or ''}")
        except Exception as e:
            db_conn.rollback()
            st.error(f"Error loading activity log: {e}")

APP_BUILD = "build-2025-10-20-11-xx"
st.caption(f"App: {APP_BUILD}")